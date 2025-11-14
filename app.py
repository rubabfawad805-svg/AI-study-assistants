import streamlit as st
import os
from groq import Groq
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import tempfile
import warnings
warnings.filterwarnings("ignore", message="Thread 'MainThread': missing ScriptRunContext!")
# -------------------------------------------------------
#  Initialize GROQ Client
# -------------------------------------------------------
def init_groq():
    api_key ="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"
    if not api_key:
        st.error("Missing GROQ_API_KEY environment variable.")
        return None
    return Groq(api_key=api_key)


# -------------------------------------------------------
#  LLM Query
# -------------------------------------------------------
def ask_groq(prompt, model="mixtral-8x7b-32768"):
    client = init_groq()
    if client is None:
        return "❌ Missing API key."
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are an AI Study Assistant."},
                {"role": "user", "content": prompt},
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"


# -------------------------------------------------------
#  File Handlers
# -------------------------------------------------------
def extract_text_from_pdf(file):
    with pdfplumber.open(file) as pdf:
        return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())


def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_image(file):
    img = Image.open(file)
    return pytesseract.image_to_string(img)


# -------------------------------------------------------
#  Study Tools Prompt Generators
# -------------------------------------------------------
def build_prompt(task, text, extra=None):
    if task == "Summarize":
        return f"Summarize the following text:\n\n{text}"

    if task == "Flashcards":
        return f"Create detailed study flashcards from this content:\n\n{text}"

    if task == "MCQs":
        return f"Generate 10 MCQs (4 options each) from:\n\n{text}"

    if task == "True/False":
        return f"Generate 10 True/False questions from:\n\n{text}"

    if task == "Short Questions":
        return f"Generate 10 short-answer questions from:\n\n{text}"

    if task == "Explain Simply":
        return f"Explain this content simply:\n\n{text}"

    if task == "Key Points":
        return f"Extract the key points:\n\n{text}"

    if task == "Study Plan":
        days = extra if extra else 7
        return f"Create a {days}-day study schedule for:\n\n{text}"

    if task == "Ask a Question":
        return f"Answer this question based on the study notes:\n\n{text}"

    return "Invalid task."


# -------------------------------------------------------
#  Streamlit App
# -------------------------------------------------------
st.set_page_config(page_title="AI Study Assistant", layout="wide")
st.title("📚 AI Study Assistant (GROQ Powered)")

st.sidebar.header("Controls")

task = st.sidebar.selectbox(
    "Choose a task:",
    [
        "Summarize",
        "Flashcards",
        "MCQs",
        "True/False",
        "Short Questions",
        "Explain Simply",
        "Key Points",
        "Study Plan",
        "Ask a Question"
    ]
)

study_plan_days = None
if task == "Study Plan":
    study_plan_days = st.sidebar.number_input("Days", min_value=1, max_value=60, value=7)

uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX, IMAGE):", type=["pdf", "docx", "pptx", "png", "jpg"])

manual_text = st.text_area("Or paste text here:")

process_button = st.button("Generate")

# -------------------------------------------------------
#  MAIN LOGIC
# -------------------------------------------------------
if process_button:
    if uploaded_file:
        file_type = uploaded_file.name.split(".")[-1].lower()

        if file_type == "pdf":
            extracted_text = extract_text_from_pdf(uploaded_file)
        elif file_type == "docx":
            extracted_text = extract_text_from_docx(uploaded_file)
        elif file_type == "pptx":
            extracted_text = extract_text_from_pptx(uploaded_file)
        elif file_type in ["png", "jpg"]:
            extracted_text = extract_text_from_image(uploaded_file)
        else:
            st.error("Unsupported file type.")
            extracted_text = ""
    else:
        extracted_text = manual_text

    if not extracted_text.strip():
        st.error("No text found. Upload a file or enter text.")
    else:
        prompt = build_prompt(task, extracted_text, study_plan_days)
        st.subheader("🧠 AI Output")
        output = ask_groq(prompt)
        st.write(output)
